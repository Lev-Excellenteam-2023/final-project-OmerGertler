import asyncio
import os
import time

import aiohttp
from flask import json
from pptx import Presentation
import openai


REQ = "Please provide an English summary and brief explanation of the text on this slide, " \
      "along with a recommended website for further learning. " \
      "If the slide does not contain any text, please do not provide an answer. " \
      "Please format your response as an educational paragraph.\n"
KEY = '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'

messages = []
messages.append({"role": "system", "content": f"you are a helpful AI that is explaining a PowerPoint presentation."})
openai.api_key = KEY


def import_to_file(text_list):
    """
    Importing list of strings of GPT answers into txt file
    :param text_list: list of strings
    """
    try:
        with open('res.txt', 'w') as file:
            [file.write(f'{line}\n') for line in text_list]
        print('\nfile is ready')
    except Exception as ex:
        print(ex)


def extract_text_from_slides(presentation_path):
    """
    Extract text from each slide from pptx file.
    :param presentation_path: The path of tje pptx file.
    :return: list of strings, each string contains slide's text.
    """
    ppt = Presentation(presentation_path)
    text_list = []
    for slide in ppt.slides:
        slide_text = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        text_list.append(''.join(slide_text))
    return text_list


def create_request(text):
    """
    Create request for GPT from the "text" parameter and return GPT's answer.
    :param text: The text from the slide.
    :return: The answer of GPT about the text parameter.
    """
    request = REQ + text
    return get_ans_from_AI_async(request)


async def get_ans_from_AI_async(user_message: str):
    """
    Send a message to GPT and return its reply.
    :param msg: msg to GPT.
    :return: The answer of GPT.
    """
    async with aiohttp.ClientSession() as session:
        response = await session.post(
            "https://api.openai.com/v1/chat/completions",  # Updated endpoint URL
            headers={
                "Authorization": f"Bearer {openai.api_key}",
                "Content-Type": "application/json"
            },
            json={
                "messages": [{"role": "system", "content": "You are a helpful assistant."},
                             {"role": "user", "content": user_message}],
                "max_tokens": 512,
                "model": "gpt-3.5-turbo"
            }
        )
        return await response.json()


async def response_handler(slides):
    """
    Handles the requests of explaining the slides to chat-gpt.

    :param slides: A list of slide contents.
    :return: A list of responses from chat-gpt for each slide.
    """
    tasks = []
    for slide_content in slides:
        task = asyncio.create_task(create_request(slide_content))
        tasks.append(task)

    responses = await asyncio.gather(*tasks)
    return responses


def save_as_json(filename, outputs_folder, result):
    output_filename = os.path.splitext(filename)[0] + '.json'
    output_path = os.path.join(outputs_folder, output_filename)
    with open(output_path, 'w') as output_file:
        json.dump(result, output_file)
    return output_filename


def process_files():
    uploads_folder = os.path.join(os.getcwd(), 'uploads')
    outputs_folder = os.path.join(os.getcwd(), 'outputs')

    while True:
        time.sleep(10)  # Sleep for 10 seconds between iterations

        # Scan the uploads folder and identify files that haven't been processed
        for filename in os.listdir(uploads_folder):
            file_path = os.path.join(uploads_folder, filename)
            if os.path.isfile(file_path) and filename not in os.listdir(outputs_folder):
                print(f"Processing file: {filename}")

                # Process the file using existing code
                data_list = extract_text_from_slides(file_path)
                loop = asyncio.get_event_loop()
                responses = loop.run_until_complete(response_handler(data_list))
                result = []
                for i, response in enumerate(responses, start=1):
                    answer = response["choices"][0]["message"]["content"]
                    result.append({'number: ': i, 'slide answer: ': answer})

                # Save the explanation JSON in the outputs folder
                output_filename = save_as_json(filename, outputs_folder, result)

                print(f"Explanation JSON saved: {output_filename}")
                os.remove(file_path)  # Remove the processed file from the uploads folder

            else:
                print(f"Skipping file: {filename} (already processed)")


if __name__ == "__main__":
    process_files()


